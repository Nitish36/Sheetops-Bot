import os
import secrets
import smartsheet
from flask import Flask, render_template, request, jsonify, redirect, url_for, session, flash
from flask_login import LoginManager, UserMixin, login_user, logout_user, login_required, current_user
from authlib.integrations.flask_client import OAuth
from google import genai
from google.genai import types
from dotenv import load_dotenv
import re
import csv
import json
import io
import boto3
from psycopg2.extras import RealDictCursor
import PyPDF2
import docx
from pptx import Presentation
import pandas as pd
from flask_wtf.csrf import CSRFProtect, generate_csrf
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
from cryptography.fernet import Fernet
from db_manager import log_activity, create_session, save_message, get_user_sessions, get_session_messages, create_user, verify_user_login, get_user_by_id, get_db_connection, get_executive_stats
from gsheet_manager import gsheet_sync_user, gsheet_log_session, gsheet_log_activity, gsheet_log_feedback
from werkzeug.middleware.proxy_fix import ProxyFix
from crawlers.general_announcement_crawler import get_community_announcements
from crawlers.events_crawler import get_smartsheet_events
from crawlers.product_announcement_crawler import get_product_updates
from crawlers.pmo_crawler import get_pmo_trends
from crawlers.healthcare_ls_crawler import get_healthcare_trends
from crawlers.financial_services_crawler import get_finance_trends
from crawlers.digital_it import get_it_trends
from crawlers.best_practices_crawler import get_best_practices
from crawlers.b2b_crawler import get_b2b_trends
from crawlers.ai_crawler import get_ai_trends
from crawlers.unanswered import get_unanswered_questions
from rag_manager import rag

load_dotenv()

app = Flask(__name__,template_folder='template')
app.wsgi_app = ProxyFix(app.wsgi_app, x_for=1, x_proto=1, x_host=1, x_prefix=1)
app.secret_key = os.getenv("SECRET_KEY")
# Config
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
#SMARTSHEET_ACCESS_TOKEN = os.getenv("SMARTSHEET_ACCESS_TOKEN")
SUPPORT_SHEET_ID = os.getenv("SUPPORT_SHEET_ID")
csrf = CSRFProtect(app)
ENCRYPTION_KEY = os.getenv("ENCRYPTION_KEY")
cipher = Fernet(ENCRYPTION_KEY.encode()) if ENCRYPTION_KEY else None
app.config.update(
    SESSION_COOKIE_SECURE=True,    # Only send over HTTPS
    SESSION_COOKIE_HTTPONLY=True,  # Prevent JS access to session cookie
    SESSION_COOKIE_SAMESITE='Lax', # Prevent CSRF
)
#CURRENT_USER_EMAIL = "nitish.pkv@gmail.com"
#CURRENT_USER_NAME = "Nitish K"

client = genai.Client(api_key=GOOGLE_API_KEY)
#sm_client = smartsheet.Smartsheet(SMARTSHEET_ACCESS_TOKEN)

# Initialize Login Manager
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = "login"

# Initialize OAuth
oauth = OAuth(app)
google = oauth.register(
    name='google',
    client_id=os.getenv("GOOGLE_CLIENT_ID"),
    client_secret=os.getenv("GOOGLE_CLIENT_SECRET"),
    server_metadata_url='https://accounts.google.com/.well-known/openid-configuration',
    client_kwargs={'scope': 'openid email profile'}
)

limiter = Limiter(
    get_remote_address,
    app=app,
    default_limits=["200 per day", "50 per hour"],
    storage_uri="memory://",
)

SYSTEM_PROMPT = """[SYSTEM: You are the SheetOps AI Assistant for Smartsheet
Admins.

TASK HIERARCHY:
1. DATA AUDIT: If you receive [SHEET DATA ATTACHED], perform a deep analysis.
Do NOT open a support ticket for this. Just provide the report.
2. INTERNAL KNOWLEDGE: You MUST prioritize the rules found in the <KNOWLEDGE_BASE> section below.
3. MANDATORY TAG: If (and only if) the user agrees to a ticket, include
TRIGGER_TICKET] in your response.]
4. READ-ONLY: You provide analysis and suggestions. You NEVER attempt to write to sheets.
5. ACTIONABLE STEPS: When you find an error (overdue tasks, security risks, formula breaks),
provide a 'Manual Remediation Plan' with specific steps the Admin should take in the Smartsheet UI.
6. COMPLIANCE: Remind the user of standard Smartsheet best practices (e.g., using 'Contact List' columns
for owners instead of text strings).]\n\n

[KNOWLEDGE MANAGEMENT RULES]:
1. You have access to multiple internal documents labeled with <DOCUMENT_START>.
2. If the user asks about Formulas or Syntax, prioritize info within 'formulas.md'.
3. If the user asks about Naming or Permissions, prioritize 'governance.md'.
4. If a user says to guide him/her on the app, prioritze 'onboarding.md'.
5. When you cite a rule, mention the specific document name (e.g., 'Per our Formula Standards...').
6. If the user asks about Premium Apps (Dynamic View, Data Shuttle, DataMesh, Bridge, Pivot App, WorkApps), prioritize 'asl.md'.
7. If the user asks about external integrations (Jira, Salesforce, ServiceNow, Teams, Slack), prioritize 'connectors.md'.
8. If the user reports an ERROR, a BROKEN SYNC, or asks for a HEALTH CHECK, you MUST prioritize 'troubleshooting.md'.
9. LIVE DATA: You have a "Live Crawler" tool.
    For general updates, use 'announcements'.
    For feature/software updates, use 'product releases'.
    For training, use 'events'.
    For PMO governance, use 'PMO trends'.
    For Healthcare, use 'Healthcare trends'.
    For Finance, use 'Financial trends'.
    For IT, Software Development, or Digital Transformation, use 'IT trends'.
    For optimization and standards, use 'best practices'.
    For vendor management or client-facing operations, use 'B2B trends'.
    Future Tech: For Smartsheet AI, Gemini, or Automation features, use 'AI trends'.
    Community Contribution: For unanswered questions or users wanting to help others, use 'unanswered questions'.
    
[ONBOARDING & GUIDANCE]: 
If a user is new, asks "How do I start?", or asks about SheetOps features, prioritize 'onboarding_guide.md'. 
Act as a mentor to help them navigate the 16 modules.

[ADVANCED LEARNING & PREMIUM FEATURES]:
When a user asks about scaling their solution or handling large datasets, reference 'asl.md' (Advanced Smartsheet Learning). 
Explain how tools like Data Shuttle or DataTable solve limitations like row counts or manual data entry. 
Maintain a strategic consultant tone for these topics.

[LIVE COMMUNITY ACCESS & UPDATES]:
When users ask about the latest Smartsheet changes or product releases, acknowledge that you can perform live crawls of the Smartsheet Community. 
If the user asks a vague question about news, guide them: "I can scrape the latest announcements for you—just ask for 'top 10 community updates'."

[LIVE DATA CAPABILITIES]:
Mention that you can pull live event schedules (Webinars, ENGAGE tours, and Training) directly from the official Smartsheet events portal.

[COMMUNICATION STYLE]:
1. Maintain your authoritative 'Lead Architect' tone. 
2. Always cite the specific internal document you are referencing (e.g., 'According to our Onboarding Guide...').
3. If a user asks a general question, always try to link it back to a specific tool in our Admin Toolkit (e.g., 
   'To solve this, I recommend using the Webhook Audit module in your Toolkit').

[ERROR RESOLUTION & HEALTH CHECKS]:
When a user describes a failure (e.g., "The sync stopped"), consult 'troubleshooting.md'. 
Check for common culprits: Column Formulas, Service Account passwords, or missing Error Columns. 
Always ask the user: "Do you see an error message in your Error Column?"

[CHART INSTRUCTION]: If your analysis contains numerical status counts (e.g., Overdue vs. On-Track), 
you MUST append a JSON block at the very end of your response in this EXACT format:
[CHART_DATA: {"labels": ["Label1", "Label2"], "values": [10, 5], "colors": ["#14b8a6", "#ef4444"]}]
Use #14b8a6 for positive/complete and #ef4444 for overdue/errors.\n\n
[TICKET TRIGGER]: If the user wants to raise a support ticket, you MUST output ONLY the phrase: 
INITIATE_HC_TICKET_FLOW
Do not ask any questions yourself. The system will take over.
"""

FORMULA_EXPERT_PROMPT = """[SYSTEM: You are the Smartsheet Formula Genius.
Your goal is to debug or write formulas with 100% accurate Smartsheet syntax.
RULES:
1. Always use [Column Name]@row for row-level references.
2. If a column name has spaces, use brackets [ ].
3. Explain WHY a formula failed (e.g., #UNPARSEABLE usually means a missing
comma or mismatched parenthesis).
4. For cross-sheet references, use curly braces {Reference Name}.
5. Format the final formula clearly in `code blocks`.]\n\n"""

USER_AUDIT_PROMPT = """[SYSTEM: You are the SheetOps User & Connector Auditor.
Analyze the 'User Management' export. Your mission is to find 'Hidden Bloat'.

ANALYSIS FOCUS:
1. CONNECTOR OVER-PROVISIONING: Identify users with 'JIRA' or 'Salesforce'
Connector Admin rights who haven't logged in recently. These are security
risks.
2. PREMIUM APP UTILIZATION: Flag users with access to 'Dynamic View',
'DataMesh', or 'Data Shuttle' who have 'Sheets (Created)' = 0. They may not
need these premium tools.
3. TRIAL CLEANUP: List users 'On Trial' and check their 'Trial End Date'.
Recommend if they should be converted or removed.
4. INACTIVE LICENSES: Licensed Users with no 'Last Login (UTC)' in 90 days.

Provide a 'Security & Premium App Summary' with an Action Plan.]\n\n"""

SEAT_AUDIT_PROMPT = """[SYSTEM: You are the SheetOps Strategic License Auditor.
You are analyzing a 'Seat Management' export from Smartsheet.

CRITICAL METRICS TO ANALYZE:
1. GHOST LICENSES: Identify 'Seat Type: Member' users who have 0 '# of Edits'
and 0 '# of Created Assets' in the last 90 days.
2. TRUE-UP RISKS: Flag users 'Qualified For True Up' who have no 'Last Paid
Activity' in the last 180 days.
3. ACTIVITY TRENDS: Compare '# of Asset Views' vs '# of Edits'. If a user only
'Views' but never 'Edits', suggest downgrading them to a Free Viewer to save
costs.
4. ADMIN OVERLOAD: Identify how many 'System Admins' and 'Group Admins' we
have. Is it too many for our org size?
5. ROI SUMMARY: Provide a table of 'Top 10 Cost Saving Opportunities' showing
User, Seat Type, and why they should be downgraded.]\n\n"""

HYGIENE_AUDIT_PROMPT = """[SYSTEM: You are the SheetOps Asset Hygiene
Specialist. Analyze the provided Sheet Inventory list.

YOUR MISSION:
1. ZOMBIE ASSETS: Identify sheets where 'createdAt' and 'modifiedAt' are
nearly identical and older than 90 days (Created but never used).
2. ABANDONED SHEETS: Identify sheets that haven't been modified in over 180
days.
3. ACCESS RISK: Summarize the 'accessLevel' distribution. Flag sheets where
the user is 'OWNER' vs 'ADMIN'.
4. NAMING GOVERNANCE: Check if sheet names look professional or if they are
'Untitled' or 'Test' sheets.

Provide a 'Clean-up Priority List' and estimate how much 'clutter' can be
removed.]\n\n"""

REPORT_AUDIT_PROMPT = """[SYSTEM: You are the SheetOps Report Strategy
Specialist. Analyze the provided Smartsheet Report Inventory.

YOUR MISSION:
1. REPORT TYPE ANALYSIS: Distinguish between 'Summary Reports'
(isSummaryReport: True) and 'Row Reports'. Explain the balance of your
portfolio.
2. ACCESS RISK: Identify reports where the user has 'OWNER' vs 'ADMIN' vs
'EDITOR' access. Flag high-access reports for security review.
3. NAMING GOVERNANCE: Identify reports with vague names (e.g., 'Copy of...',
'Test', 'Untitled').
4. PERMALINK AUDIT: Verify if reports are using the standard Smartsheet domain
and look for any unusual URL patterns.

Provide a 'Report Governance Summary' and suggest a better naming convention
for the 'Untitled' or 'Test' assets found.]\n\n"""

DASHBOARD_AUDIT_PROMPT = """[SYSTEM: You are the SheetOps Dashboard Experience
Specialist. Analyze the provided Smartsheet Dashboard Inventory.

YOUR MISSION:
1. EXECUTIVE CLUTTER: Identify dashboards that haven't been modified in over
180 days. These are likely 'Dead Dashboards' that provide no current value.
2. ZOMBIE DASHBOARDS: Identify dashboards created and never modified
(createdAt == modifiedAt).
3. VISIBILITY AUDIT: Summarize the 'accessLevel' distribution. Who has
'OWNER' rights to our executive views?
4. NAMING CONVENTIONS: Look for unprofessional names like 'Test', 'My
Dashboard', or 'New Dashboard'.

Provide a 'Visibility Action Plan' to clean up our reporting layer.]\n\n"""

WEBHOOK_AUDIT_PROMPT = """[SYSTEM: You are the SheetOps Webhook Security
Specialist.Analyze the provided Smartsheet Webhook Inventory.

YOUR MISSION:
1. INTEGRATION HEALTH: Identify webhooks where 'enabled' is False. Deep-dive
into 'status' and 'disabledDetails' (e.g., CALLBACK_FAILED, SCOPE_INACCESSIBLE)
to explain why it's broken.
2. SECURITY SCAN: Analyze 'callbackUrl'. Flag any non-HTTPS URLs or suspicious/
unknown domains.
3. PERFORMANCE AUDIT: Check 'stats' (lastSuccessfulCallback vs
lastCallbackAttempt). Flag webhooks that haven't succeeded in over 30 days.
4. APP MONITORING: Group by 'apiClientName'. Identify if any third-party apps
are dominating the webhook count.
5. ZOMBIE CLEANUP: Identify webhooks that are 'NEW_NOT_VERIFIED' and older
than 7 days—these are abandoned setups.

Provide a 'Webhook Security & Health Report' with a priority list of broken
integrations to fix.]\n\n"""

WORKSPACE_AUDIT_PROMPT = """[SYSTEM: You are the SheetOps Workspace Governance
Specialist.Analyze the provided Smartsheet Workspace Inventory.

YOUR MISSION:
1. ACCESS DISTRIBUTION: Identify how many workspaces the user is an 'OWNER' of
versus 'ADMIN' or 'EDITOR'.
2. CONTAINER SPRAWL: Identify if there are too many 'Personal' or 'Test'
workspaces that should be consolidated into Departmental workspaces.
3. SECURITY REVIEW: Highlight any workspaces with 'OWNER' access that might
need to be transferred (e.g., if naming suggests it belongs to a different
department).
4. NAMING GOVERNANCE: Flag workspaces with vague names like 'My Workspace',
'New Workspace', or 'Stuff'.

Provide a 'Workspace Governance Report' with recommendations for better
container organization.]\n\n"""

PUBLISHED_AUDIT_PROMPT = """[SYSTEM: You are the SheetOps Data Privacy &
Security Auditor. Analyze the provided Smartsheet Published Items list.

YOUR MISSION:
1. SECURITY RISK ASSESSMENT: Identify items where 'Access Control' is set
to 'Anyone' or 'Public'. Flag these as high-priority security reviews.
2. STALE PUBLIC DATA: Identify published items that haven't been modified
in over 120 days ('Last Modified Date/Time'). Suggest unpublishing if the
data is no longer current.
3. PUBLISHER AUDIT: Group by 'Publisher'. Identify if specific users are
publishing a high volume of assets.
4. FORMAT ANALYSIS: Check 'Published Format' (e.g., Read Only, Edit, iCal).
Flag any 'Edit' access links as extreme security risks.

Provide a 'Public Exposure & Security Report' with a list of links that
should be restricted or disabled immediately.]\n\n"""

WORKAPPS_AUDIT_PROMPT = """[SYSTEM: You are the SheetOps WorkApps Utilization
Auditor.Analyze the provided WorkApps Collaborator Report.

YOUR MISSION:
1. COLLABORATOR EFFICIENCY: Identify users who haven't accessed a WorkApp
('Last App Access Date/Time') in over 90 days. Recommend removing them from
the Collaborator Pack to save costs.
2. EXTERNAL RISK: Identify 'External' collaborators. Flag any external users
who have access to sensitive-sounding apps.
3. OWNER AUDIT: Group by 'App Owner'. Identify if any apps are owned by users
who might no longer be in the relevant department.
4. USAGE DENSITY: Look at 'Total # of WorkApps accessed'. Identify 'Power
Users' vs 'Low Users'.
5. BILLING IMPACT: Analyze the 'Collab Pack Purchased' column. Are we
over-utilizing or under-utilizing our purchased packs?

Provide a 'WorkApps ROI & Security Report' with a list of users to remove
to optimize the Collaborator Pack.]\n\n"""

ATTACHMENT_AUDIT_PROMPT = """[SYSTEM: You are the SheetOps Storage & Security Auditor.
Analyze the provided Smartsheet Attachment Inventory.

YOUR MISSION:
1. STORAGE BLOAT: Identify the largest files (sizeInKb). Flag files over 10MB (10240Kb) as candidates for external storage (OneDrive/SharePoint).
2. SECURITY SCAN: Check 'attachmentType' and 'mimeType'. Flag potentially dangerous files (e.g., .exe, .zip, .msi) or non-standard types.
3. STALE DATA: Identify attachments created ('createdAt') over 1 year ago on sheets that might no longer be active.
4. USER UPLOAD TRENDS: Group by 'createdBy'. Identify which users are consuming the most storage quota.
5. RECOMMENDATION: Provide a 'Storage Optimization Plan' with estimated space that can be reclaimed.]\n\n"""

PDF_SUMMARIZER_PROMPT = """[SYSTEM: You are the SheetOps Technical Document Analyst.
Your goal is to summarize complex Smartsheet or Corporate Governance documents.

ANALYSIS REQUIREMENTS:
1. EXECUTIVE SUMMARY: Provide a 3-sentence high-level overview.
2. KEY TAKEAWAYS: Extract the top 5 most important rules or technical requirements.
3. ADMIN ACTIONS: List specific actions a Smartsheet Admin must take based on this document.
4. JARGON CLEANUP: Explain any complex terms in simple language.]\n\n"""


DOCX_SUMMARIZER_PROMPT = """[SYSTEM: You are the SheetOps Document Intelligence Expert.
Your goal is to summarize Microsoft Word documents containing Smartsheet governance, project plans, or technical requirements.

ANALYSIS REQUIREMENTS:
1. EXECUTIVE SUMMARY: High-level overview of the document purpose.
2. ACTION ITEMS: Bullet points of specific tasks for a Smartsheet Admin.
3. LOGIC & RULES: Extract any specific formulas, naming conventions, or workflow rules mentioned.
4. TIMELINES: Identify key dates or milestones if present.]\n\n"""

PPTX_SUMMARIZER_PROMPT = """[SYSTEM: You are the SheetOps Presentation Strategist.
Your goal is to analyze PowerPoint decks related to Smartsheet projects, status reports, or training.

ANALYSIS REQUIREMENTS:
1. DECK OVERVIEW: Provide a summary of the presentation's goal.
2. KEY MILESTONES: Identify any project dates, phases, or deadlines mentioned across slides.
3. SMARTSHEET IMPACT: Highlight specific mentions of sheet structures, reports, or automation needs.
4. Q&A PREP: Suggest 3 questions an Admin should ask the presenter based on these slides.]\n\n"""

EXCEL_ANALYST_PROMPT = """[SYSTEM: You are the SheetOps Data Architect.
Your goal is to analyze random Excel files and provide a structural and strategic summary.

ANALYSIS REQUIREMENTS:
1. DATA STRUCTURE: List the columns found and identify the 'Primary Key' or main identifier.
2. CONTENT SUMMARY: Describe what this data represents (e.g., Budget, Task List, Contact Sheet).
3. SMARTSHEET MAPPING: Suggest which Smartsheet Column Types should be used for this data (e.g., 'Date' column for Col A, 'Contact List' for Col B).
4. DATA ANOMALIES: Flag any missing values or inconsistent formatting you detect in the sample.]\n\n"""

DASHBOARD_BUILDER_PROMPT = """[SYSTEM: You are the SheetOps AI Data Scientist. 
Analyze the provided data summary and design a professional executive dashboard.

OUTPUT REQUIREMENTS (Strict JSON format only):
Return a JSON object with these keys:
1. "metrics": An array of 4 objects { "label": "String", "value": "Number/String", "color": "Teal/Red/Blue/Amber" }
2. "charts": An array of 2 objects { "type": "bar/doughnut", "title": "String", "labels": [], "values": [], "colors": [] }
3. "summary": A 3-sentence high-level insight about the data.

Use professional colors: Teal (#14b8a6), Red (#ef4444), Blue (#3b82f6), Amber (#f59e0b).]\n\n"""

TICKET_OPTIONS = {
    "REQUEST_TYPES": ["Workspace Creation", "Create Blueprint", "Other"],
    "WS_TYPES": ["Digital PMO Workspace Creation","Workspace Creation for MDM Implementation","Workday Workspace Creation",
        "Healthcare Engagement Health", "Others"],
    "CC_L1": ["Digital", "Education", "Healthcare", "Human Resources", "Business Advisory Consulting", "Corporate Finance Consulting", "Corporate", "Global Enablement", "Information Technology", "Innosight", "Legal", "Marketing", "Operations"],
    "CC_L2": ["Consulting-BA-Financial Advisory","Consulting-BA-Transformation","Consulting-CF-Financial Institutions Advisory",
  "Consulting-CF-Healthcare FAS","Consulting-EDR-Global Philanthropy","Consulting-EDR-RS-Research Services","Consulting-EDR-RS-Sales & Admin",
  "Consulting-EDR-RS-Software Services","Consulting-EDR-Strategy & Operations","Consulting-HC Studer Education","Consulting-HC-Assess and Design","Consulting-HC-PI Implementations","Consulting-HC-PI Insights",
  "Consulting-HC-People Implementations","Consulting-HC-Sustained Performance Services","Corporate-Executive","Corporate-Facilities","Corporate-Finance & Accounting","Corporate-Global Enablement","Corporate-Global Finance",
  "Corporate-Global HR Operations","Corporate-HR-Business Partners","Corporate-HR-Delivery Operations","Corporate-HR-Learning & Talent","Corporate-HR-Total Rewards",
  "Corporate-IT-Applications","Corporate-IT-Infrastructure & Service Center","Corporate-IT-Product Services","Corporate-Marketing","Corporate-Operations","Corporate-Recruiting","Digital-Admin",
  "Digital-Business Intelligence and Data Solutions","Digital-CRM","Digital-Data Management","Digital-Digital Workforce Enablement","Digital-EDR-RS-Revenue & Delivery","Digital-EDR-Spend Management",
  "Digital-Emerging Technology-Tagetik","Digital-Engagement Assurance","Digital-HC-HMS Software","Digital-HC-HTSG Services","Digital-HC-People Software","Digital-HC-Software Product Implementation","Digital-India-Business Intelligence and Data Solutions",
  "Digital-India-Data Management","Digital-India-Oracle","Digital-India-Program Management Office","Digital-Intelligent Automation","Digital-Managed Services",
  "Digital-Markets-Admin","Digital-Markets-Global Analytics","Digital-Oracle","Digital-Product-Admin","Digital-Products-Customer Success","Digital-Products-Engineering-Apps",
  "Digital-Products-HC Sales","Digital-Products-Product Management","Digital-Project Management Office","Digital-Salesforce","Digital-Technology Market Leaders","Digital-Workday","EDR-Admin","EDR-Analytics","HC-Admin",
  "HC-Portfolio & Strategy Management","HC-Sales Operations","Managed Services-HC Admin","Managed Services-HC-Adventist","Managed Services-HC-Offshore"],

    "SERVICE_LINE": ["Corporate", "Engagement Health", "HRS", "Oracle", "Other", "Other Service Line", "Salesforce", "Workday"],
    "CATEGORY": ["Access Request","Archiving closed Projects (including following up and verifying readiness to Archive for one Project)","Automations","Blueprint","Charts","Dashboards","Data Shuttle","Data Mesh","Discussion/Meeting",
  "Dynamic View","Forms","Full Solution","Global Roll outs","Metrics KPIs","Other","PIVOT","Project Workspace Creation (automation/Manual)","Reports","Requirement Analysis for new Program (End to End Process)","Sheets",
  "Vacation Calendar setup","Workapps","Group Management"]

}
WS_FORM_LINKS = {
    "Digital PMO Workspace Creation": "https://app.smartsheet.com/b/form/digital-pmo-example",
    "Workspace Creation for MDM Implementation": "https://app.smartsheet.com/b/form/mdm-example",
    "Workday Workspace Creation": "https://app.smartsheet.com/b/form/workday-example",
    "Healthcare Engagement Health": "https://app.smartsheet.com/b/form/healthcare-example",
    "Others": "https://app.smartsheet.com/b/form/general-workspace-example"
}

# User Class for Flask-Login
class User(UserMixin):
    def __init__(self, id, email, name, plan, username):
        self.id = id
        self.email = email
        self.username = username
        self.name = name if name else username  # Fallback if name is null
        self.plan = plan


def get_sm_client():
    """Fetches and decrypts token from DB instead of Session."""
    if not current_user.is_authenticated:
        return None

    conn = get_db_connection()
    cur = conn.cursor(cursor_factory=RealDictCursor)
    cur.execute("SELECT encrypted_token FROM users WHERE id = %s", (current_user.id,))
    res = cur.fetchone()
    cur.close()
    conn.close()

    if res and res['encrypted_token']:
        try:
            decrypted_token = cipher.decrypt(res['encrypted_token'].encode()).decode()
            return smartsheet.Smartsheet(decrypted_token)
        except Exception as e:
            print(f"Decryption error: {e}")
    return None

@app.route('/check-connection')
@login_required
def check_connection():
    client = get_sm_client() # This now checks the DB
    return jsonify({"connected": client is not None})

@app.route('/save-smartsheet-token', methods=['POST'])
@login_required
def save_sm_token():
    data = request.json
    token = data.get('token')
    if token and cipher:
        # Encrypt the token before saving
        encrypted = cipher.encrypt(token.encode()).decode()

        conn = get_db_connection()
        cur = conn.cursor()
        cur.execute("UPDATE users SET encrypted_token = %s WHERE id = %s", (encrypted, current_user.id))
        conn.commit()
        cur.close()
        conn.close()
        return jsonify({"status": "success"})
    return jsonify({"status": "error"}), 400

@app.after_request
def set_csrf_cookie(response):
    response.set_cookie('csrf_token', generate_csrf(), secure=True, samesite='Lax')
    return response

@login_manager.user_loader
def load_user(user_id):
    u = get_user_by_id(user_id) # This function is in db_manager.py
    if u:
        # Use .get() to avoid errors if the key is missing
        return User(
            id=u['id'],
            email=u['email'],
            username=u['username'],
            name=u.get('name', u['username']),
            plan=u.get('plan', 'free')
        )
    return None


@app.route('/login/google')
def login_google():
    nonce = secrets.token_urlsafe(16)
    session['nonce'] = nonce
    redirect_uri = url_for('google_callback', _external=True)
    return google.authorize_redirect(redirect_uri, nonce=nonce)


@app.route('/login/callback')
def google_callback():
    nonce = session.pop('nonce', None)
    if not nonce:
        flash("Login expired. Please try again.", "error")
        return redirect('/login')

    try:
        token = google.authorize_access_token()
        user_info = google.parse_id_token(token, nonce=nonce)

    except Exception as e:
        print(f"Auth Security Error: {e}")
        flash("Security verification failed.", "error")
        return redirect('/login')

    email = user_info['email']
    name = user_info['name']
    username = user_info['name']

    # --- HURON VERIFICATION LOGIC ---
    # Check if email ends with @hcg (or the full corporate domain like @hcg.com)
    # Adjust '@hcg' to your exact corporate domain suffix
    plan_type = 'free'
    if email.lower().endswith('@hcg.com'):
        plan_type = 'huron'

    # Save/Update User in Neon DB
    user_id = ensure_user_exists_with_plan(email, name, plan_type)

    # Create the user object and log them in
    user_obj = User(user_id, email,username, name, plan_type)
    login_user(user_obj)
    gsheet_sync_user(user_id, email, name, plan_type)

    return redirect('/')# Back to dashboard


def ensure_user_exists_with_plan(email, name, plan):
    conn = get_db_connection()
    cur = conn.cursor(cursor_factory=RealDictCursor)
    try:
        cur.execute("SELECT id FROM users WHERE email = %s", (email,))
        user = cur.fetchone()
        if user:
            cur.execute("UPDATE users SET last_login = CURRENT_TIMESTAMP, plan = %s WHERE id = %s", (plan, user['id']))
            user_id = user['id']
        else:
            cur.execute("INSERT INTO users (email, name, plan) VALUES (%s, %s, %s) RETURNING id", (email, name, plan))
            user_id = cur.fetchone()['id']
        conn.commit()
        return user_id
    finally:
        cur.close()
        conn.close()


# --- ROUTE 1: USER & CONNECTOR AUDIT ---
@app.route('/user-audit', methods=['POST'])
@login_required
def user_audit():
    file = request.files.get('file')
    log_activity(current_user.id, "User Audit", "upload_csv")

    if not file:
        return jsonify({"response": "No file."})

    try:
        stream = io.StringIO(file.stream.read().decode("UTF8"), newline=None)
        csv_reader = csv.DictReader(stream)
        rows = []
        for i, row in enumerate(csv_reader):
            if i > 100:
                break
            rows.append({
                "User": row.get('Email'),
                "Licensed": row.get('Licensed User'),
                "Jira": row.get('JIRA Connector Admin'),
                "Salesforce": row.get('Salesforce Connector Admin'),
                "DynView": row.get('Dynamic View'),
                "DataMesh": row.get('DataMesh'),
                "Shuttle": row.get('Data Shuttle'),
                "SheetsCreated": row.get('Sheets (Created)'),
                "LastLogin": row.get('Last Login (UTC)'),
                "OnTrial": row.get('On Trial')
            })

        response = client.models.generate_content(
            model='gemini-3-flash-preview',
            contents=[{"role": "user", "parts": [{"text": USER_AUDIT_PROMPT + f"DATA:\n{str(rows)}"}]}]
        )
        gsheet_log_activity(current_user.name, "User Audit", "CSV Upload")
        return jsonify({"response": response.text})
    except Exception as e:
        return jsonify({"response": f"Error: {str(e)}"})


# --- ROUTE 2: Seat  AUDIT ---
@app.route('/seat-audit', methods=['POST'])
@login_required
def seat_audit():
    if 'file' not in request.files:
        return jsonify({"response": "No file uploaded."})

    file = request.files['file']
    log_activity(current_user.id, "Seat Audit", "upload_csv")

    try:
        stream = io.StringIO(file.stream.read().decode("UTF8"), newline=None)
        csv_reader = csv.DictReader(stream)

        rows = []
        for i, row in enumerate(csv_reader):
            if i > 100:
                break  # Keep it fast for Vercel

            # Map the exact columns you provided
            rows.append({
                "Email": row.get('Email'),
                "Status": row.get('User Status'),
                "Seat": row.get('Seat Type'),
                "TrueUp": row.get('Qualified For True Up'),
                "SysAdmin": row.get('System Admin'),
                "Edits90": row.get('# of Edits (Last 90 days)'),
                "Created365": row.get('# of Created Assets (Last 365 days)'),
                "LastPaidAct": row.get('Last Paid Activity (UTC)'),
                "LastLogin": row.get('Last Login (UTC)')
            })

        contents = [{
            "role": "user",
            "parts": [{"text": SEAT_AUDIT_PROMPT + f"DATASET:\n{str(rows)}"}]
        }]

        response = client.models.generate_content(
            model='gemini-3-flash-preview',
            contents=contents
        )
        gsheet_log_activity(current_user.name, "Seat Audit", "CSV Upload")
        return jsonify({"response": response.text})
    except Exception as e:
        return jsonify({"response": f"Seat Audit Error: {str(e)}"})


# Route 3 Sheet Audit
@app.route('/hygiene-audit', methods=['POST'])
@login_required
def hygiene_audit():
    file = request.files.get('file')
    log_activity(current_user.id, "Sheet Audit", "upload_csv")
    if not file:
        return jsonify({"response": "No file."})

    try:
        stream = io.StringIO(file.stream.read().decode("UTF8"), newline=None)
        csv_reader = csv.DictReader(stream)
        rows = []
        for i, row in enumerate(csv_reader):
            if i > 200:
                break  # Sheet lists are simple, we can take more rows
            rows.append({
                "Name": row.get('name'),
                "Access": row.get('accessLevel'),
                "Created": row.get('createdAt'),
                "Modified": row.get('modifiedAt')
            })

        response = client.models.generate_content(
            model='gemini-3-flash-preview', 
            contents=[{"role": "user", "parts": [{"text": HYGIENE_AUDIT_PROMPT + f"DATA:\n{str(rows)}"}]}]
        )
        gsheet_log_activity(current_user.name, "Sheet Audit", "CSV Upload")
        return jsonify({"response": response.text})
    except Exception as e:
        return jsonify({"response": f"Error: {str(e)}"})


# Route 4 Report Audit
@app.route('/report-audit', methods=['POST'])
@login_required
def report_audit():
    file = request.files.get('file')
    log_activity(current_user.id, "Report Audit", "upload_csv")
    if not file:
        return jsonify({"response": "No file."})

    try:
        stream = io.StringIO(file.stream.read().decode("UTF8"), newline=None)
        csv_reader = csv.DictReader(stream)
        rows = []
        for i, row in enumerate(csv_reader):
            if i > 200:
                break

            # UPDATED MAPPING: Matching your new columns
            rows.append({
                "ID": row.get('id'),
                "Name": row.get('name'),
                "Access": row.get('accesslevel'),
                "Link": row.get('permalink'),
                "IsSummary": row.get('isSummaryReport')
            })

        response = client.models.generate_content(
            model='gemini-3-flash-preview', 
            contents=[{"role": "user", "parts": [{"text": REPORT_AUDIT_PROMPT + f"DATA:\n{str(rows)}"}]}]
        )
        gsheet_log_activity(current_user.name, "Report Audit", "CSV Upload")
        return jsonify({"response": response.text})
    except Exception as e:
        return jsonify({"response": f"Error: {str(e)}"})


# Route 5 Dashboard Audit
@app.route('/dashboard-audit', methods=['POST'])
@login_required
def dashboard_audit():
    file = request.files.get('file')
    log_activity(current_user.id, "Dashboard Audit", "upload_csv")
    if not file:
        return jsonify({"response": "No file."})

    try:
        stream = io.StringIO(file.stream.read().decode("UTF8"), newline=None)
        csv_reader = csv.DictReader(stream)
        rows = []
        for i, row in enumerate(csv_reader):
            if i > 200:
                break
            rows.append({
                "Name": row.get('name'),
                "Access": row.get('accessLevel'),
                "Created": row.get('createdAt'),
                "Modified": row.get('modifiedAt')
            })

        response = client.models.generate_content(
            model='gemini-3-flash-preview', 
            contents=[{"role": "user", "parts": [{"text": DASHBOARD_AUDIT_PROMPT + f"DATA:\n{str(rows)}"}]}]
        )
        gsheet_log_activity(current_user.name, "Dashboard Audit", "CSV Upload")
        return jsonify({"response": response.text})
    except Exception as e:
        return jsonify({"response": f"Error: {str(e)}"})


# Route 6 Webhooks Audit
@app.route('/webhook-audit', methods=['POST'])
@login_required
def webhook_audit():
    file = request.files.get('file')
    log_activity(current_user.id, "Webhook Audit", "upload_csv")
    if not file:
        return jsonify({"response": "No file."})

    try:
        stream = io.StringIO(file.stream.read().decode("UTF8"), newline=None)
        csv_reader = csv.DictReader(stream)
        rows = []
        for i, row in enumerate(csv_reader):
            if i > 150:
                break

            # Mapping your exact columns
            rows.append({
                "Name": row.get('name'),
                "URL": row.get('callbackUrl'),
                "Status": row.get('status'),
                "Enabled": row.get('enabled'),
                "App": row.get('apiClientName'),
                "Stats": row.get('stats'),
                "DisabledDetails": row.get('disabledDetails'),
                "Modified": row.get('modifiedAt')
            })

        response = client.models.generate_content(
            model='gemini-3-flash-preview', 
            contents=[{"role": "user", "parts": [{"text": WEBHOOK_AUDIT_PROMPT + f"DATA:\n{str(rows)}"}]}]
        )
        gsheet_log_activity(current_user.name, "Webhook Audit", "CSV Upload")
        return jsonify({"response": response.text})
    except Exception as e:
        return jsonify({"response": f"Error: {str(e)}"})

# Route 7 Workspace Audit
@app.route('/workspace-audit', methods=['POST'])
@login_required
def workspace_audit():
    file = request.files.get('file')
    log_activity(current_user.id, "Workspace Audit", "upload_csv")
    if not file:
        return jsonify({"response": "No file."})

    try:
        stream = io.StringIO(file.stream.read().decode("UTF8"), newline=None)
        csv_reader = csv.DictReader(stream)
        rows = []
        for i, row in enumerate(csv_reader):
            if i > 200:
                break

            # Mapping your exact columns: id, name, accessLevel, permalink
            rows.append({
                "ID": row.get('id'),
                "Name": row.get('name'),
                "Access": row.get('accessLevel'),
                "Link": row.get('permalink')
            })

        response = client.models.generate_content(
            model='gemini-3-flash-preview', 
            contents=[{"role": "user", "parts": [{"text": WORKSPACE_AUDIT_PROMPT + f"DATA:\n{str(rows)}"}]}]
        )
        gsheet_log_activity(current_user.name, "Workspace Audit", "CSV Upload")
        return jsonify({"response": response.text})
    except Exception as e:
        return jsonify({"response": f"Error: {str(e)}"})


# Route 8 Published Audit
@app.route('/published-audit', methods=['POST'])
@login_required
def published_audit():
    file = request.files.get('file')
    log_activity(current_user.id, "Publish Audit", "upload_csv")
    if not file:
        return jsonify({"response": "No file."})

    try:
        stream = io.StringIO(file.stream.read().decode("UTF8"), newline=None)
        csv_reader = csv.DictReader(stream)
        rows = []
        for i, row in enumerate(csv_reader):
            if i > 150:
                break

            # Mapping your exact columns
            rows.append({
                "Name": row.get('Name'),
                "Type": row.get('Type'),
                "Publisher": row.get('Publisher'),
                "Format": row.get('Published Format'),
                "Access": row.get('Access Control'),
                "Link": row.get('Publish Link'),
                "Modified": row.get('Last Modified Date/Time (UTC)')
            })

        response = client.models.generate_content(
            model='gemini-3-flash-preview', 
            contents=[{"role": "user", "parts": [{"text": PUBLISHED_AUDIT_PROMPT + f"DATA:\n{str(rows)}"}]}]
        )
        # 3. Use the ACTUAL logged-in username for Google Sheet/Looker logs
        gsheet_log_activity(current_user.name, "Published Audit", "CSV Upload")
        return jsonify({"response": response.text})
    except Exception as e:
        return jsonify({"response": f"Error: {str(e)}"})


# Route 9 Workapp Audit
@app.route('/workapps-audit', methods=['POST'])
@login_required  # 1. Ensure security
def workapps_audit():
    file = request.files.get('file')

    # 2. Use the ACTUAL logged-in user ID for Postgres logs
    log_activity(current_user.id, "Workapp Audit", "upload_csv")

    if not file:
        return jsonify({"response": "No file."})

    try:
        stream = io.StringIO(file.stream.read().decode("UTF8"), newline=None)
        csv_reader = csv.DictReader(stream)
        rows = []
        for i, row in enumerate(csv_reader):
            if i > 150:
                break

            rows.append({
                "CollabType": row.get('Internal/External Paid Collaborator'),
                "User": row.get('WorkApps Collaborator'),
                "AppName": row.get('Last App Accessed'),
                "Owner": row.get('App Owner'),
                "LastAccess": row.get('Last App Access Date/Time (UTC)'),
                "TotalAccessed": row.get('Total # of WorkApps accessed that month'),
                "PackType": row.get('Collab Pack Purchased at the time')
            })

        response = client.models.generate_content(
            model='gemini-3-flash-preview',
            contents=[{"role": "user", "parts": [{"text": WORKAPPS_AUDIT_PROMPT + f"DATA:\n{str(rows)}"}]}]
        )

        # 3. Use the ACTUAL logged-in username for Google Sheet/Looker logs
        gsheet_log_activity(current_user.name, "Workapp Audit", "CSV Upload")

        return jsonify({"response": response.text})
    except Exception as e:
        print(f"WorkApps Audit Error: {e}")
        return jsonify({"response": f"Error: {str(e)}"})


# Route 10 Log Tool Click
@app.route('/log-tool-click', methods=['POST'])
@login_required  # Ensure user is logged in
def log_tool_click():
    data = request.json
    tool_name = data.get("tool_name")

    # Use ACTUAL logged-in user info
    log_activity(current_user.id, tool_name, "click")
    gsheet_log_activity(current_user.name, tool_name, "Toolkit Click")

    return jsonify({"status": "logged"})


# Route 11 Register Route
@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        username = request.form.get('username')
        email = request.form.get('email')
        password = request.form.get('password')

        result = create_user(username, email, password)
        if result['success']:
            user_obj = User(result['user_id'], email, username,username, result['plan'])
            login_user(user_obj)
            gsheet_sync_user(result['user_id'], email, username, result['plan'])
            flash("Welcome to SheetOps! Account created successfully.", "success")
            return redirect('/')
        else:
            return render_template('register.html', error=result['error'])

    return render_template('register.html')


# Route 12 login Route
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        email = request.form.get('email')
        password = request.form.get('password')

        user_data = verify_user_login(email, password)
        if user_data:
            user_obj = User(
                id=user_data['id'],
                email=user_data['email'],
                username=user_data['username'],
                name=user_data.get('name', user_data['username']),  # Fallback to username if name is missing
                plan=user_data.get('plan', 'free')
            )
            login_user(user_obj)
            gsheet_sync_user(user_obj.id, user_obj.email, user_obj.name, user_obj.plan)
            flash(f"Welcome back, {user_data['username']}!", "success")
            return redirect('/')
        else:
            flash("Invalid credentials. Please try again.", "error")
            return render_template('login.html', error="Invalid credentials. Try again.")

    return render_template('login.html')

# Route 13 PDF Summarizer
@app.route('/summarize-pdf', methods=['POST'])
@login_required
def summarize_pdf():
    if 'file' not in request.files:
        return jsonify({"response": "No file uploaded."})

    file = request.files['file']

    try:
        # 1. Extract Text from PDF
        pdf_reader = PyPDF2.PdfReader(file)
        text_content = ""
        # Limit to first 10 pages to keep it fast for Vercel
        for page_num in range(min(len(pdf_reader.pages), 10)):
            text_content += pdf_reader.pages[page_num].extract_text()

        # 2. Log Activity
        log_activity(current_user.id, "PDF Summarizer", "upload_pdf")
        gsheet_log_activity(current_user.name, "PDF Summarizer", "File Uploaded")

        # 3. Send to Gemini 3
        contents = [{
            "role": "user",
            "parts": [{"text": PDF_SUMMARIZER_PROMPT + f"DOCUMENT TEXT:\n{text_content}"}]
        }]

        response = client.models.generate_content(
            model='gemini-3-flash-preview',
            contents=contents
        )

        return jsonify({"response": response.text})

    except Exception as e:
        print(f"PDF Error: {e}")
        return jsonify({"response": f"Failed to process PDF: {str(e)}"})


# Route 14 Attachment Route
@app.route('/attachment-audit', methods=['POST'])
@login_required
def attachment_audit():
    file = request.files.get('file')
    if not file: return jsonify({"response": "No file uploaded."})

    # Log Activity
    log_activity(current_user.id, "Attachment Audit", "upload_csv")
    gsheet_log_activity(current_user.name, "Attachment Audit", "CSV Uploaded")

    try:
        stream = io.StringIO(file.stream.read().decode("UTF8"), newline=None)
        csv_reader = csv.DictReader(stream)
        rows = []
        for i, row in enumerate(csv_reader):
            if i > 100: break  # Keep performance snappy
            rows.append({
                "Sheet": row.get('sheetName'),
                "FileName": row.get('attachmentName'),
                "Type": row.get('attachmentType'),
                "SizeKB": row.get('sizeInKb'),
                "Mime": row.get('mimeType'),
                "Created": row.get('createdAt'),
                "User": row.get('createdBy')
            })

        response = client.models.generate_content(
            model='gemini-3-flash-preview',
            contents=[{"role": "user", "parts": [{"text": ATTACHMENT_AUDIT_PROMPT + f"DATA:\n{str(rows)}"}]}]
        )
        return jsonify({"response": response.text})
    except Exception as e:
        return jsonify({"response": f"Error: {str(e)}"})


# Route 15 Word Docx Summarizer
@app.route('/summarize-docx', methods=['POST'])
@login_required
def summarize_docx():
    file = request.files.get('file')
    if not file: return jsonify({"response": "No file uploaded."})

    try:
        # 1. Extract Text from DOCX
        doc = docx.Document(file)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)

        text_content = "\n".join(full_text)

        # 2. Log Activity
        log_activity(current_user.id, "Word Summarizer", "upload_docx")
        gsheet_log_activity(current_user.name, "Word Summarizer", "File Uploaded")

        # 3. Send to Gemini 3 (AI Client should be available)
        contents = [{
            "role": "user", 
            "parts": [{"text": DOCX_SUMMARIZER_PROMPT + f"DOCUMENT TEXT:\n{text_content[:15000]}"}] # Capped for speed
        }]

        response = client.models.generate_content(
            model='gemini-3-flash-preview', 
            contents=contents
        )

        return jsonify({"response": response.text})

    except Exception as e:
        print(f"Word Doc Error: {e}")
        return jsonify({"response": f"Failed to process Word Doc: {str(e)}"})

# Route 16 Powerpoint Summarizer
@app.route('/summarize-pptx', methods=['POST'])
@login_required
def summarize_pptx():
    file = request.files.get('file')
    if not file: return jsonify({"response": "No file uploaded."})
    
    try:
        # 1. Extract Text from PPTX
        prs = Presentation(file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        
        text_content = "\n".join(text_runs)

        # 2. Log Activity
        log_activity(current_user.id, "PPTX Summarizer", "upload_pptx")
        gsheet_log_activity(current_user.name, "PPTX Summarizer", "File Uploaded")

        # 3. Send to Gemini 3
        contents = [{
            "role": "user", 
            "parts": [{"text": PPTX_SUMMARIZER_PROMPT + f"PRESENTATION TEXT:\n{text_content[:15000]}"}]
        }]

        response = client.models.generate_content(
            model='gemini-3-flash-preview', 
            contents=contents
        )

        return jsonify({"response": response.text})

    except Exception as e:
        print(f"PPTX Error: {e}")
        return jsonify({"response": f"Failed to process PowerPoint: {str(e)}"})

#Route 17 Excel analyst
@app.route('/summarize-excel', methods=['POST'])
@login_required
def summarize_excel():
    file = request.files.get('file')
    if not file: return jsonify({"response": "No file uploaded."})
    
    try:
        # 1. Read Excel using Pandas
        df = pd.read_excel(file)
        
        # 2. Create a text-based summary of the data
        # We send the column names and the first 15 rows to Gemini
        data_summary = f"Columns: {list(df.columns)}\n\nSample Data:\n{df.head(15).to_string()}"

        # 3. Log Activity
        log_activity(current_user.id, "Excel Analyst", "upload_excel")
        gsheet_log_activity(current_user.name, "Excel Analyst", "File Uploaded")

        # 4. Send to Gemini 3
        contents = [{
            "role": "user", 
            "parts": [{"text": EXCEL_ANALYST_PROMPT + f"EXCEL DATA SAMPLE:\n{data_summary}"}]
        }]

        response = client.models.generate_content(
            model='gemini-3-flash-preview', 
            contents=contents
        )

        return jsonify({"response": response.text})

    except Exception as e:
        print(f"Excel Error: {e}")
        return jsonify({"response": f"Failed to process Excel file: {str(e)}"})

# Route 18 Dashboard Builder
@app.route('/generate-dashboard', methods=['POST'])
@login_required
def generate_dashboard():
    file = request.files.get('file')
    if not file: return jsonify({"error": "No file"}), 400

    try:
        # 1. Read file safely
        filename = file.filename.lower()
        if filename.endswith('.csv'):
            df = pd.read_csv(file)
        elif filename.endswith(('.xlsx', '.xls')):
            # Ensure openpyxl is in requirements.txt
            df = pd.read_excel(file, engine='openpyxl')
        else:
            return jsonify({"error": "Unsupported file format"}), 400

        # 2. Create a compact summary for the AI
        # We only send the column names and some basic stats to stay within limits
        data_summary = {
            "columns": list(df.columns),
            "rows_count": len(df),
            "sample": df.head(5).to_dict(orient='records')
        }

        # 3. Call Gemini 3
        # Note: We pass the JSON instruction in the prompt directly to be safer
        prompt = DASHBOARD_BUILDER_PROMPT + f"\nDATA SUMMARY:\n{json.dumps(data_summary)}"

        response = client.models.generate_content(
            model='gemini-3-flash-preview',
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type='application/json'
            )
        )

        # 4. Cleanup AI response (Gemini sometimes adds backticks)
        raw_json = response.text.strip()
        if raw_json.startswith("```json"):
            raw_json = raw_json.replace("```json", "").replace("```", "").strip()

        # Parse it back to Python to verify, then send to JS
        dashboard_data = json.loads(raw_json)

        log_activity(current_user.id, "Dashboard Builder", "generate_dashboard")
        return jsonify(dashboard_data)

    except Exception as e:
        print(f"DASHBOARD BUILDER ERROR: {str(e)}")
        import traceback
        traceback.print_exc()  # This will show the real error in Vercel Logs
        return jsonify({"error": f"Backend Error: {str(e)}"}), 500

@app.route('/logout')
@login_required
def logout():
    logout_user()
    flash(f"Logged Out Successfully!", "success")
    return redirect('/')


def handle_ticketing_flow(user_input, session_id):
    if 't_data' not in session:
        session['t_data'] = {}

    # --- STEP 0: INITIAL TRIGGER ---
    if user_input == "INITIATE_HC_TICKET_FLOW":
        session['t_state'] = 0
        return jsonify({
            "response": "I'll help you raise a corporate support ticket. First, confirm or enter your **Corporate Email**:",
            "session_id": str(session_id)
        })

    state = session.get('t_state')

    # --- STEP 1: SELECT REQUEST TYPE ---
    if state == 0:
        session['t_data']['email'] = user_input
        session['t_state'] = 1
        return jsonify({
            "response": "Thank you. Select the **type of request**:",
            "options": TICKET_OPTIONS["REQUEST_TYPES"],
            "session_id": str(session_id)
        })

    # --- STEP 2: BRANCHING LOGIC ---
    elif state == 1:
        session['t_data']['request_type'] = user_input

        if user_input == "Create Blueprint":
            session['t_state'] = 20  # Jump to Blueprint Branch
            return jsonify({
                "response": "Blueprint Request initiated. What is your **Cost Center Level 1**?",
                "options": TICKET_OPTIONS["CC_L1"],
                "session_id": str(session_id)
            })
        elif user_input == "Workspace Creation":
            session['t_state'] = 10  # Jump to Workspace Branch
            return jsonify({
                "response": "Workspace Request initiated. What is your **Cost Center Level 1**?",
                "options": TICKET_OPTIONS["CC_L1"],
                "session_id": str(session_id)
            })
        else:
            session['t_state'] = 2  # Jump to "Other" Branch
            return jsonify({
                "response": "Support request initiated. What is your **Cost Center Level 1**?",
                "options": TICKET_OPTIONS["CC_L1"],
                "session_id": str(session_id)
            })

    # ==========================================
    # BRANCH: CREATE BLUEPRINT (States 20-23)
    # ==========================================
    elif state == 20:  # CC L1 received
        session['t_data']['cc_l1'] = user_input
        session['t_state'] = 21
        return jsonify({
            "response": f"Selected {user_input}. Step 2: Select **Cost Center Level 2**:",
            "options": TICKET_OPTIONS["CC_L2"],
            "session_id": str(session_id)
        })

    elif state == 21:  # CC L2 received
        session['t_data']['cc_l2'] = user_input
        session['t_state'] = 22
        return jsonify({
            "response": "Step 3: Please enter the **Leader/Director Name** for this blueprint:",
            "session_id": str(session_id)
        })

    elif state == 22:  # Leader Name received
        session['t_data']['leader_name'] = user_input
        session['t_state'] = 23
        return jsonify({
            "response": "Step 4: Almost done. **Please provide the details for Blueprint Creation** (Goals, required sheets, etc.):",
            "session_id": str(session_id)
        })

    elif state == 23:  # Blueprint Details received -> FINAL SAVE
        session['t_data']['description'] = user_input
        success = log_hcg_ticket_to_smartsheet(session['t_data'])

        session.pop('t_state', None)
        session.pop('t_data', None)

        if success:
            return jsonify({
                "response": "✅ **Blueprint Request Logged.** Our team will review the details and reach out for a kickoff session.",
                "session_id": str(session_id)
            })
        return jsonify({"response": "⚠️ Error writing to Smartsheet.", "session_id": str(session_id)})

    # ==========================================
    # BRANCH: WORKSPACE CREATION (States 10-12)
    # ==========================================
    elif state == 10:
        session['t_data']['cc_l1'] = user_input
        session['t_state'] = 11
        return jsonify({"response": "Select **Cost Center Level 2**:", "options": TICKET_OPTIONS["CC_L2"],
                        "session_id": str(session_id)})

    elif state == 11:
        session['t_data']['cc_l2'] = user_input
        session['t_state'] = 12
        return jsonify({"response": "Select the **Workspace Creation Type**:", "options": TICKET_OPTIONS["WS_TYPES"],
                        "session_id": str(session_id)})

    elif state == 12:
        ws_type = user_input
        session['t_data']['ws_type'] = ws_type
        session['t_data']['description'] = f"Workspace creation request for type: {ws_type}"
        success = log_hcg_ticket_to_smartsheet(session['t_data'])
        form_link = WS_FORM_LINKS.get(ws_type, WS_FORM_LINKS["Others"])

        session.pop('t_state', None)
        session.pop('t_data', None)
        return jsonify({
                           "response": f"✅ **Workspace Request Logged.** Please complete the official form: [Click Here]({form_link})",
                           "session_id": str(session_id)})

    # ==========================================
    # BRANCH: OTHER (States 2-6)
    # ==========================================
    elif state == 2:
        session['t_data']['cc_l1'] = user_input
        session['t_state'] = 3
        return jsonify({"response": "Select **Cost Center Level 2**:", "options": TICKET_OPTIONS["CC_L2"],
                        "session_id": str(session_id)})

    elif state == 3:
        session['t_data']['cc_l2'] = user_input
        session['t_state'] = 4
        return jsonify({"response": "Select the **Service Line**:", "options": TICKET_OPTIONS["SERVICE_LINE"],
                        "session_id": str(session_id)})

    elif state == 4:
        session['t_data']['service_line'] = user_input
        session['t_state'] = 5
        return jsonify({"response": "What is the **Category** of this request?", "options": TICKET_OPTIONS["CATEGORY"],
                        "session_id": str(session_id)})

    elif state == 5:
        session['t_data']['category'] = user_input
        session['t_state'] = 6
        return jsonify({"response": "Please provide a brief **Task Description**.", "session_id": str(session_id)})

    elif state == 6:
        session['t_data']['description'] = user_input
        success = log_hcg_ticket_to_smartsheet(session['t_data'])
        session.pop('t_state', None)
        session.pop('t_data', None)
        return jsonify({"response": "✅ **Support Ticket Logged.**", "session_id": str(session_id)})

    return jsonify({"response": "I encountered an error in the ticket flow. Let's restart."})


def log_hcg_ticket_to_smartsheet(data):
    """Writes the gathered form data to your Smartsheet Support Sheet."""
    try:
        sm_personal = get_sm_client()
        if not sm_personal: return False

        sheet = sm_personal.Sheets.get_sheet(SUPPORT_SHEET_ID)
        new_row = smartsheet.models.Row()
        new_row.to_top = True

        # This matches the column titles in your sheet
        fields_map = {
            "Requestor's Email": data.get('email'),
            "Request Type": data.get('request_type'),
            "Cost Center Level 1": data.get('cc_l1'),
            "Cost Center Level 2": data.get('cc_l2'),
            "Workspace Type": data.get('ws_type', 'N/A'),
            "Leader/Director Name": data.get('leader_name', 'N/A'),
            "Service Line": data.get('service_line', 'N/A'),
            "Category": data.get('category', 'N/A'),
            "Task Description": data.get('description', '')
        }

        for col in sheet.columns:
            if col.title in fields_map:
                new_row.cells.append({'column_id': col.id, 'value': str(fields_map[col.title])})

        sm_personal.Sheets.add_rows(SUPPORT_SHEET_ID, [new_row])
        return True
    except Exception as e:
        print(f"Smartsheet Log Error: {e}")
        return False



@app.route('/get-sessions', methods=['GET'])
@login_required  # Add this to ensure only logged-in users can fetch history
def get_sessions():
    # Use current_user.id directly (No need to call ensure_user_exists here)
    sessions = get_user_sessions(current_user.id)

    # Convert UUIDs and Datetimes to strings for JSON
    for s in sessions:
        s['id'] = str(s['id'])
        s['created_at'] = s['created_at'].strftime("%Y-%m-%d %H:%M")

    return jsonify(sessions)

@app.route('/get-stats', methods=['GET'])
@login_required
def get_stats():
    stats = get_executive_stats(current_user.id)
    return jsonify(stats)


@app.route('/get-session-chat/<session_id>', methods=['GET'])
@login_required
def get_session_chat(session_id):
    messages = get_session_messages(session_id)
    return jsonify(messages)


def get_sheet_data_for_audit(sheet_id, sm_client):
    """Fetches a sheet using the provided personal sm_client."""
    try:
        # Now using the client passed into the function
        sheet = sm_client.Sheets.get_sheet(sheet_id)

        cols = [col.title for col in sheet.columns]
        rows_summary = []
        for row in sheet.rows:
            row_dict = {}
            for i, cell in enumerate(row.cells):
                col_name = cols[i]
                row_dict[col_name] = cell.display_value
            rows_summary.append(row_dict)

        return {
            "name": sheet.name,
            "columns": cols,
            "data": rows_summary[:100]
        }
    except Exception as e:
        # Catching all errors (404, 403, or invalid token)
        return {"error": f"Connection Error: {str(e)}"}

"""def load_knowledge_base():
    kb_content = ""
    kb_folder = "knowledge"
    if os.path.exists(kb_folder):
        for filename in os.listdir(kb_folder):
            if filename.endswith(".md"):
                # We add clear markers so the AI knows which 'Law Book' it is reading
                kb_content += f"\n<DOCUMENT_START: {filename}>\n"
                with open(os.path.join(kb_folder, filename), "r", encoding="utf-8") as f:
                    kb_content += f.read()
                kb_content += f"\n<DOCUMENT_END: {filename}>\n"
    return kb_content

# Initialize the knowledge once when the server starts
INTERNAL_KNOWLEDGE = load_knowledge_base()"""


@app.route('/formula', methods=['POST'])
@login_required
def formula_builder():
    data = request.json
    user_query = data.get("query")

    try:
        # We don't need history here; every formula is a fresh start
        contents = [{
            "role": "user",
            "parts": [{"text": FORMULA_EXPERT_PROMPT + user_query}]
        }]

        response = client.models.generate_content(
            model='gemini-3-flash-preview',
            contents=contents
        )

        return jsonify({"response": response.text})
    except Exception as e:
        return jsonify({"response": f"Formula Analysis Error: {str(e)}"})


@app.route('/')
def index():
    return render_template('index.html', user=current_user)


@app.route('/log-feedback', methods=['POST'])
@login_required
def log_feedback():
    data = request.json
    fb_type = data.get("feedback_type")
    response_id = data.get("response_id")

    # 1. Log to Postgres
    log_activity(current_user.id, "AI Response", f"Feedback: {fb_type}")

    # 2. Log to Google Sheets for Looker
    gsheet_log_feedback(current_user.name, fb_type, response_id)

    return jsonify({"status": "logged"})


@app.route('/admin/reindex', methods=['GET'])
@login_required  # Highly recommended so strangers can't trigger it
def reindex_knowledge():
    try:
        # If you have an admin flag on your User model, check it here:
        # if not current_user.is_admin: return "Unauthorized", 403

        # This clears the old memory and re-scans the /knowledge folder
        rag.load_and_index()

        return jsonify({
            "status": "success",
            "message": "Knowledge base re-indexed successfully!",
            "total_chunks": len(rag.embeddings),
            "files_processed": list(set(chunk['source'] for chunk in rag.chunks))
        })
    except Exception as e:
        return jsonify({
            "status": "error",
            "message": f"Failed to re-index: {str(e)}"
        }), 500


@app.route('/chat', methods=['POST'])
@login_required
@limiter.limit("5 per minute")
def chat():
    data = request.json
    user_message = data.get("message")
    raw_history = data.get("history", [])

    # 1. SANITIZE SESSION ID (Prevents "null" string errors)
    session_id = data.get("session_id")
    if session_id == "null" or not session_id:
        session_id = None

    # --- NEW: TICKETING STATE MACHINE INTERCEPTION ---
    # Check if a ticket flow is currently active in this session
    t_state = session.get('t_state')

    # If the user wants to start a ticket OR is already in the middle of one
    if user_message == "INITIATE_HC_TICKET_FLOW" or t_state is not None:
        return handle_ticketing_flow(user_message, session_id)

    # --- NEW: COMMUNITY ANNOUNCEMENT INTERCEPTION ---
    # Detects keywords like 'announcement', 'community news', 'whats new'
    if any(word in user_message.lower() for word in ["announcement", "community news", "what's new"]):
        # Extract number if user specified (e.g., "Give me 5 announcements")
        num_match = re.search(r'\d+', user_message)
        limit = int(num_match.group()) if num_match else 10

        # This triggers the "Progress Bar" state on the frontend
        announcements = get_community_announcements(limit)

        # Save this interaction to DB before returning
        if not session_id:
            session_id = create_session(current_user.id, "Community Scrape", "chats")
        save_message(session_id, "user", user_message)
        save_message(session_id, "model", f"Scraped {len(announcements)} announcements.")

        return jsonify({
            "type": "announcements",
            "data": announcements,
            "response": f"I've crawled the Smartsheet Community for the latest updates. Here are the top {len(announcements)} announcements:",
            "session_id": str(session_id)
        })

    # --- EVENT INTERCEPTION ---
    if any(word in user_message.lower() for word in ["event", "webinar", "training", "conference"]):
        num_match = re.search(r'\d+', user_message)
        limit = int(num_match.group()) if num_match else 5  # Events are long, default to 5

        events = get_smartsheet_events(limit)

        # Save to DB logic (Reuse your existing save_message logic)
        return jsonify({
            "type": "events",
            "data": events,
            "response": f"I've found {len(events)} upcoming Smartsheet events and webinars for you:",
            "session_id": str(session_id)
        })

    # 1. Check for Product Specific News first
    product_keywords = ["product release", "new feature", "product announcement", "updates"]
    if any(word in user_message.lower() for word in product_keywords):
        num_match = re.search(r'\d+', user_message)
        limit = int(num_match.group()) if num_match else 10
        updates = get_product_updates(limit)

        return jsonify({
            "type": "product_news",
            "data": updates,
            "response": f"I've pulled the latest {len(updates)} Product Announcements and feature releases for you:",
            "session_id": str(session_id)
        })

    # Check for PMO / Trending topics
    pmo_keywords = ["pmo", "project management office", "pmo trends", "pmo advice", "pmo best practices"]
    if any(word in user_message.lower() for word in pmo_keywords):
        num_match = re.search(r'\d+', user_message)
        limit = int(num_match.group()) if num_match else 10
        trends = get_pmo_trends(limit)

        return jsonify({
            "type": "pmo_trends",
            "data": trends,
            "response": f"I've analyzed the latest trending PMO discussions for you. Here are the top {len(trends)} 'Hot' topics:",
            "session_id": str(session_id)
        })

    # Check for Healthcare / Life Sciences topics
    hc_keywords = ["healthcare", "life sciences", "medical", "hipaa", "clinical", "pharma", "patient"]

    if any(word in user_message.lower() for word in hc_keywords):
        num_match = re.search(r'\d+', user_message)
        limit = int(num_match.group()) if num_match else 10
        trends = get_healthcare_trends(limit)

        return jsonify({
            "type": "healthcare_trends",
            "data": trends,
            "response": f"I've pulled the latest trending Healthcare & Life Sciences discussions for you. Here are the top {len(trends)} topics:",
            "session_id": str(session_id)
        })

    # Check for Financial Services / Banking topics
    finance_keywords = ["finance", "financial", "banking", "insurance", "investment", "accounting", "audit"]
    if any(word in user_message.lower() for word in finance_keywords):
        num_match = re.search(r'\d+', user_message)
        limit = int(num_match.group()) if num_match else 10
        trends = get_finance_trends(limit)

        return jsonify({
            "type": "finance_trends",
            "data": trends,
            "response": f"I've retrieved the latest trending Financial Services discussions for you. Here are the top {len(trends)} topics:",
            "session_id": str(session_id)
        })

    # Check for IT / Digital Transformation topics
    it_keywords = ["digital", " it ", "tech", "software", "sdlc", "infrastructure", "digital transformation", "it pmo"]
    if any(word in user_message.lower() for word in it_keywords):
        num_match = re.search(r'\d+', user_message)
        limit = int(num_match.group()) if num_match else 10
        trends = get_it_trends(limit)

        return jsonify({
            "type": "it_trends",
            "data": trends,
            "response": f"I've analyzed the latest trending Digital IT & Portfolio discussions. Here are the top {len(trends)} topics:",
            "session_id": str(session_id)
        })

    # Check for Best Practices / Optimization topics
    bp_keywords = ["best practice", "optimization", "efficiency", "how to improve", "expert advice", "standards"]
    if any(word in user_message.lower() for word in bp_keywords):
        num_match = re.search(r'\d+', user_message)
        limit = int(num_match.group()) if num_match else 10
        practices = get_best_practices(limit)

        return jsonify({
            "type": "best_practices",
            "data": practices,
            "response": f"I've curated the latest expert Best Practices from the community. Here are the top {len(practices)} optimization topics:",
            "session_id": str(session_id)
        })

    # Check for B2B / Enterprise Work Management topics
    b2b_keywords = ["b2b", "vendor management", "client project", "external collaboration", "enterprise work"]
    if any(word in user_message.lower() for word in b2b_keywords):
        num_match = re.search(r'\d+', user_message)
        limit = int(num_match.group()) if num_match else 10
        trends = get_b2b_trends(limit)

        return jsonify({
            "type": "b2b_trends",
            "data": trends,
            "response": f"I've pulled the trending B2B Work Management discussions. Here are the top {len(trends)} topics for enterprise collaboration:",
            "session_id": str(session_id)
        })

    # Check for AI / Artificial Intelligence / Gemini topics
    ai_keywords = [" ai ", "artificial intelligence", "gemini", "generative ai", "ai features", "automation"]
    if any(word in user_message.lower() for word in ai_keywords):
        num_match = re.search(r'\d+', user_message)
        limit = int(num_match.group()) if num_match else 10
        trends = get_ai_trends(limit)

        return jsonify({
            "type": "ai_trends",
            "data": trends,
            "response": f"I've crawled the latest discussions on AI and Innovation. here are the top {len(trends)} trending topics:",
            "session_id": str(session_id)
        })

    # Check for Unanswered / Help needed topics
    help_keywords = ["unanswered", "no replies", "help others", "unresolved", "community help"]
    if any(word in user_message.lower() for word in help_keywords):
        num_match = re.search(r'\d+', user_message)
        limit = int(num_match.group()) if num_match else 10
        questions = get_unanswered_questions(limit)

        return jsonify({
            "type": "unanswered_questions",
            "data": questions,
            "response": f"I've found {len(questions)} recent discussions that haven't received an answer yet. Here's where the community needs your help:",
            "session_id": str(session_id)
        })

    # 2. STRICT HISTORY CLEANER (Ensures Gemini 3 never gets malformed data)
    clean_history = []
    for msg in raw_history:
        role = "user" if msg.get("role") == "user" else "model"
        # Safely extract text from various possible JS formats
        parts = msg.get("parts", [])
        if parts:
            if isinstance(parts[0], dict):
                txt = parts[0].get("text", "")
            else:
                txt = str(parts[0])

            if txt:
                clean_history.append({"role": role, "parts": [{"text": txt}]})


    # 3. DATABASE: Create Session & Sync User if this is a new chat
    try:
        if not session_id:
            # Create session in Postgres
            session_id = create_session(current_user.id, user_message[:30], "chats")
            # Log session to Google Sheets
            gsheet_log_session(session_id, current_user.name, user_message[:30], "chats")
            # SYNC user data to GSheet (updates last login/plan)
            gsheet_sync_user(current_user.id, current_user.email, current_user.name, current_user.plan)

        # Save the incoming User Message to DB
        save_message(session_id, "user", user_message)
    except Exception as e:
        print(f"DB Error in Chat: {e}")
        # We continue even if DB fails so the user can still chat

    # 4. PHASE 2 LOGIC: DETECT SHEET ID FOR AUDIT
    sheet_id_match = re.search(r'\b\d{15,20}\b', user_message)
    prompt_to_send = user_message

    if sheet_id_match:
        sheet_id = sheet_id_match.group()
        sm_personal = get_sm_client()

        if not sm_personal:
            prompt_to_send = "SYSTEM ERROR: No Smartsheet token found. Tell the user to go to 'API Connections' in the toolkit to provide their token before I can audit sheets."


        sheet_info = get_sheet_data_for_audit(sheet_id, sm_personal)
        if sheet_info and "error" in sheet_info:
            prompt_to_send = f"SYSTEM ALERT: The user asked to audit Sheet ID {sheet_id}, but I got this error: {sheet_info['error']}. Explain that you cannot access it."
        elif sheet_info:
            prompt_to_send = f"""
            [USER REQUEST]: {user_message}
            [DATA CONTEXT FOR AUDIT]
            Sheet Name: {sheet_info['name']}
            Columns: {sheet_info['columns']}
            Rows: {sheet_info['data']}
            [CRITICAL INSTRUCTION]: 
            1. Perform audit using ONLY the data above. 
            2. Do NOT open a ticket for this. 
            3. Do NOT include [TRIGGER_TICKET] unless explicitly asked for human help.
            4. For every issue found, provide a 'Recommended Action' (e.g., 'Update Status of Row 5 to Complete').
            5. Provide a summary 'Admin Checklist' at the end that I can copy and follow.
            6. If you find numerical status counts (e.g. Overdue, On-Track), you MUST append this tag at the very end:
                   [CHART_DATA: {{"labels": ["Overdue", "On-Track", "Complete"], "values": [5, 10, 2], "colors": ["#ef4444", "#14b8a6", "#3b82f6"]}}]
            7. Use #ef4444 for Overdue/Errors and #14b8a6 for Success/On-Track.
            """

    # 5. PREPEND SYSTEM PROMPT (Only for the very first message)
    relevant_context = rag.retrieve(user_message, top_k=3)
    if not clean_history:
        final_prompt = f"""
                    {SYSTEM_PROMPT}

                    <RELEVANT_KNOWLEDGE_CONTEXT>
                    {relevant_context}
                    </RELEVANT_KNOWLEDGE_CONTEXT>

                    USER REQUEST: {prompt_to_send}
                    """
    else:
        # For follow-up questions, we can also retrieve context or just send the message
        final_prompt = f"Context for this question: {relevant_context}\n\nUser: {prompt_to_send}"

    try:
        # 6. GEMINI API CALL
        response = client.models.generate_content(
            model='gemini-3-flash-preview',
            contents=clean_history + [{"role": "user", "parts": [{"text": final_prompt}]}]
        )

        bot_text = response.text

        # Save AI Response to DB
        save_message(session_id, "model", bot_text)

        # 7. TICKET TRIGGER LOGIC
        if "INITIATE_HC_TICKET_FLOW" in bot_text:
            return handle_ticketing_flow("INITIATE_HC_TICKET_FLOW", session_id)

        # 8. RETURN RESPONSE
        return jsonify({
            "response": bot_text,
            "session_id": str(session_id),
            "new_history_entry": [
                {"role": "user", "parts": [{"text": user_message}]},
                {"role": "model", "parts": [{"text": bot_text}]}
            ]
        })

    except Exception as e:
        print(f"CRITICAL CHAT ERROR: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({"response": f"I encountered a logic error ({str(e)[:50]}...). Please try a New Session."})


if __name__ == '__main__':
    app.run()
